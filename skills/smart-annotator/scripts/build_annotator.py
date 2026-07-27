#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
把一份内容/文件转成「可批注视图」并注入智能批注器，产出打开即用的单文件网页。

支持：md/markdown、html/htm、txt、csv（轻量，工具内解析）
     docx（mammoth→HTML）、xlsx/xls（openpyxl→表格）、pdf（逐页图片）、
     pptx（LibreOffice→pdf→逐页图片，兜底取文字）、png/jpg/gif/webp/svg（图片）

用法：
    python build_annotator.py <input_file> <output_html> [--format ...] [--title "标题"]

二进制格式（docx/xlsx/pdf/pptx/图片）会记录 originalFile，批注单里会提示：
「请对原始文件按批注修改并保留原格式」——由 Claude 用对应文档技能回写。

依赖（缺则 pip 安装）：mammoth openpyxl pymupdf python-pptx；pptx 渲染需系统 soffice(LibreOffice)。
"""
import sys, os, json, base64, argparse, subprocess, tempfile, html as _html

# ---------- 工具 ----------
def data_uri(raw_bytes, mime):
    return f"data:{mime};base64," + base64.b64encode(raw_bytes).decode()

def img_page_html(uris, note_label):
    """逐页图片视图。每页带 data-page 结构锚点，批注即可携带真实页码。"""
    parts = ['<div style="display:flex;flex-direction:column;gap:14px">']
    for i, u in enumerate(uris, 1):
        parts.append(
            f'<figure data-page="{i}" style="margin:0"><img src="{u}" alt="{note_label} 第{i}页" '
            f'style="max-width:100%;border:1px solid #ddd"/>'
            f'<figcaption style="text-align:center;color:#888;font-size:12px">第 {i} 页</figcaption></figure>')
    parts.append('</div>')
    return "\n".join(parts)

# ---------- 各格式 → (fmt, source, originalFile) ----------
def conv_docx(path, name):
    import mammoth
    with open(path, "rb") as f:
        html = mammoth.convert_to_html(f).value
    return "html", html or "<p>（空文档）</p>", {"name": name, "note": "Word 文档转 HTML 视图"}

MAX_ROWS, MAX_COLS = 500, 60   # 渲染上限，超出部分提示截断

def sheet_prefix(title):
    """Excel 引用惯例：含空格/标点的表名需用单引号包裹 —— 'My Sheet'!B3"""
    if any(ch in title for ch in " '!\"()-+*/,"):
        return "'" + title.replace("'", "''") + "'!"
    return title + "!"

def conv_xlsx(path, name):
    """Excel → 带真实单元格锚点的表格视图。
    每个单元格写入 data-cell="Sheet1!B3"，批注即可携带精确坐标，
    并渲染 Excel 式的行号/列标栏，让用户看得见自己批的是哪一格。"""
    import openpyxl
    from openpyxl.utils import get_column_letter
    wb = openpyxl.load_workbook(path, data_only=True)   # 非 read_only：需要合并单元格信息
    blocks = []
    for ws in wb.worksheets:
        n_rows, n_cols = min(ws.max_row or 0, MAX_ROWS), min(ws.max_column or 0, MAX_COLS)
        if not n_rows or not n_cols:
            continue
        pfx = sheet_prefix(ws.title)
        # 合并单元格：记录锚点跨度，以及被覆盖需跳过的格子
        span, covered = {}, set()
        for rng in ws.merged_cells.ranges:
            span[(rng.min_row, rng.min_col)] = (rng.max_row - rng.min_row + 1,
                                                rng.max_col - rng.min_col + 1)
            for r in range(rng.min_row, rng.max_row + 1):
                for c in range(rng.min_col, rng.max_col + 1):
                    if (r, c) != (rng.min_row, rng.min_col):
                        covered.add((r, c))
        # 列标栏 A B C…（gutter 不可批注）
        head = '<tr><th class="gutter"></th>' + "".join(
            f'<th class="gutter">{get_column_letter(c)}</th>' for c in range(1, n_cols + 1)) + "</tr>"
        trs = []
        for r in range(1, n_rows + 1):
            cells = [f'<th class="gutter">{r}</th>']
            for c in range(1, n_cols + 1):
                if (r, c) in covered:
                    continue
                cell = ws.cell(row=r, column=c)
                v = "" if cell.value is None else str(cell.value)
                ref = f"{pfx}{get_column_letter(c)}{r}"
                rs, cs = span.get((r, c), (1, 1))
                attr = f' rowspan="{rs}"' if rs > 1 else ""
                attr += f' colspan="{cs}"' if cs > 1 else ""
                tag = "th" if r == 1 else "td"
                cells.append(f'<{tag} data-cell="{_html.escape(ref, quote=True)}"{attr}>'
                             f'{_html.escape(v)}</{tag}>')
            trs.append("<tr>" + "".join(cells) + "</tr>")
        note = ""
        if (ws.max_row or 0) > MAX_ROWS or (ws.max_column or 0) > MAX_COLS:
            note = (f'<p class="muted">（仅渲染前 {n_rows} 行 × {n_cols} 列，'
                    f'原表 {ws.max_row}×{ws.max_column}）</p>')
        blocks.append(f"<h3>{_html.escape(ws.title)}</h3>{note}"
                      f"<table data-sheet=\"{_html.escape(ws.title, quote=True)}\">"
                      f"<thead>{head}</thead><tbody>{''.join(trs)}</tbody></table>")
    return "html", "\n".join(blocks) or "<p>（空表）</p>", {
        "name": name,
        "note": "Excel 转表格视图；批注位置为真实单元格地址（如 Sheet1!B3），请据此定位修改"}

def pdf_to_uris(path, zoom=1.6, max_pages=40):
    import fitz
    doc = fitz.open(path)
    uris = []
    for i, page in enumerate(doc):
        if i >= max_pages:
            break
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        uris.append(data_uri(pix.tobytes("png"), "image/png"))
    doc.close()
    return uris

def conv_pdf(path, name):
    uris = pdf_to_uris(path)
    return "html", img_page_html(uris, "PDF"), {"name": name, "note": "PDF 逐页图片，用圈选批注；改写由 pdf 技能处理"}

def conv_pptx(path, name):
    # 首选 LibreOffice 转 pdf 再逐页转图
    try:
        with tempfile.TemporaryDirectory() as td:
            subprocess.run(
                ["soffice", "--headless", "-env:UserInstallation=file:///tmp/lo_annot_profile",
                 "--convert-to", "pdf", "--outdir", td, path],
                check=True, capture_output=True, timeout=120)
            pdfs = [f for f in os.listdir(td) if f.lower().endswith(".pdf")]
            if pdfs:
                uris = pdf_to_uris(os.path.join(td, pdfs[0]))
                return "html", img_page_html(uris, "幻灯片"), {"name": name, "note": "PPT 逐页图片，用圈选批注；改写由 pptx 技能处理"}
    except Exception as e:
        sys.stderr.write(f"[warn] LibreOffice 渲染 pptx 失败，回退取文字：{e}\n")
    # 兜底：python-pptx 抽取每页文字
    from pptx import Presentation
    prs = Presentation(path)
    blocks = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs).strip()
                    if t:
                        texts.append(_html.escape(t))
        body = "".join(f"<p>{t}</p>" for t in texts) or "<p class='muted'>（本页无文字）</p>"
        blocks.append(f"<section style='border:1px solid #ddd;border-radius:8px;padding:12px;margin:10px 0'>"
                      f"<h3>第 {i} 页</h3>{body}</section>")
    return "html", "\n".join(blocks), {"name": name, "note": "PPT 文字视图（无渲染），按段落批注；改写由 pptx 技能处理"}

def conv_image(path, name):
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    mime = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
            "gif": "image/gif", "webp": "image/webp", "svg": "image/svg+xml"}.get(ext, "image/png")
    with open(path, "rb") as f:
        uri = data_uri(f.read(), mime)
    src = f'<div style="text-align:center"><img src="{uri}" alt="{_html.escape(name)}" style="max-width:100%"/></div>'
    return "html", src, {"name": name, "note": "图片，用圈选批注；改写需图像工具/重绘"}

def conv_text_like(path, ext):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        txt = f.read()
    if ext in ("md", "markdown"):
        return "markdown", txt, None
    if ext in ("html", "htm"):
        return "html", txt, None
    if ext == "csv":
        return "csv", txt, None
    return "text", txt, None

def convert(path, force_fmt):
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    name = os.path.basename(path)
    if force_fmt:  # 显式指定则按文本类处理
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return force_fmt, f.read(), None
    if ext == "docx":
        return conv_docx(path, name)
    if ext in ("xlsx", "xls"):
        return conv_xlsx(path, name)
    if ext == "pdf":
        return conv_pdf(path, name)
    if ext == "pptx":
        return conv_pptx(path, name)
    if ext in ("png", "jpg", "jpeg", "gif", "webp", "svg"):
        return conv_image(path, name)
    return conv_text_like(path, ext)

# ---------- 注入模板 ----------
def main():
    ap = argparse.ArgumentParser(
        description="把一个或多个文件转成可批注视图并注入智能批注器")
    ap.add_argument("paths", nargs="+",
                    help="输入文件（可多个）；未用 --out 时，最后一个视为输出 HTML（兼容旧用法）")
    ap.add_argument("--out", default=None, help="输出 HTML 路径")
    ap.add_argument("--format", choices=["html", "markdown", "csv", "text"], default=None,
                    help="强制按该文本格式处理（仅单文件时有意义）")
    ap.add_argument("--title", default=None)
    ap.add_argument("--lang", choices=["en", "zh"], default=None,
                    help="界面语言，默认英文；用户用中文交流时传 zh")
    args = ap.parse_args()

    # CLI 兼容：有 --out 则全部为输入；否则最后一个位置参数是输出
    if args.out:
        inputs, output = args.paths, args.out
    elif len(args.paths) >= 2:
        inputs, output = args.paths[:-1], args.paths[-1]
    else:
        sys.exit("缺少输出路径：用 --out out.html，或按旧用法把输出作为最后一个参数。")

    docs = []
    for path in inputs:
        fmt, source, original_file = convert(path, args.format if len(inputs) == 1 else None)
        d = {"name": os.path.basename(path), "format": fmt, "source": source}
        if original_file:
            d["originalFile"] = original_file
        docs.append(d)

    template_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                 "..", "assets", "annotator.html")
    with open(template_path, "r", encoding="utf-8") as f:
        htmldoc = f.read()

    title = args.title or (docs[0]["name"] if len(docs) == 1
                           else f"{len(docs)} documents")
    payload = {"title": title, "docs": docs}
    if args.lang:
        payload["lang"] = args.lang
    # 单文件时同时给出旧字段，保证老版本模板也能读
    if len(docs) == 1:
        payload["format"] = docs[0]["format"]
        payload["source"] = docs[0]["source"]
        if "originalFile" in docs[0]:
            payload["originalFile"] = docs[0]["originalFile"]

    js_literal = json.dumps(payload, ensure_ascii=False).replace("</", "<\\/")
    marker = "/*__PRELOAD__*/ null"
    if marker not in htmldoc:
        sys.exit("模板缺少 PRELOAD 占位符。请确认 assets/annotator.html 版本。")
    htmldoc = htmldoc.replace(marker, "/*__PRELOAD__*/ " + js_literal, 1)

    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)
    with open(output, "w", encoding="utf-8") as f:
        f.write(htmldoc)

    detail = ", ".join(f'{d["name"]}→{d["format"]}' for d in docs)
    print("已生成可批注工具：{}（{} 篇：{}；约 {} KB）".format(
        output, len(docs), detail, len(htmldoc) // 1024))

if __name__ == "__main__":
    main()
